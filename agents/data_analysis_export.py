"""Export data analysis results to PowerPoint and PDF."""

import io
import base64
from datetime import datetime


def _decode_chart_image(chart_image_b64: str | None) -> bytes | None:
    if not chart_image_b64:
        return None
    try:
        if "," in chart_image_b64:
            chart_image_b64 = chart_image_b64.split(",", 1)[1]
        return base64.b64decode(chart_image_b64)
    except Exception:
        return None


def build_pptx(data: dict, chart_image_b64: str | None = None) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    # ── Palette ──────────────────────────────────────────────────────────────
    BG      = RGBColor(0x0F, 0x17, 0x2A)   # deep navy
    BG2     = RGBColor(0x1E, 0x29, 0x3B)   # card surface
    BG3     = RGBColor(0x0B, 0x10, 0x1E)   # darker panel
    TEXT    = RGBColor(0xF1, 0xF5, 0xF9)   # off-white
    MUTED   = RGBColor(0x64, 0x74, 0x8B)   # slate-500
    SUBTLE  = RGBColor(0x94, 0xA3, 0xB8)   # slate-400

    # Section accent colours
    C = {
        "blue":   RGBColor(0x3B, 0x82, 0xF6),
        "cyan":   RGBColor(0x06, 0xB6, 0xD4),
        "green":  RGBColor(0x10, 0xB9, 0x81),
        "amber":  RGBColor(0xF5, 0x9E, 0x0B),
        "red":    RGBColor(0xEF, 0x44, 0x44),
        "purple": RGBColor(0x8B, 0x5C, 0xF6),
        "gold":   RGBColor(0xC9, 0xA9, 0x6E),   # brand Antique Gold
    }

    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Helpers ──────────────────────────────────────────────────────────────
    def _set_bg(slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = BG

    def _shape(slide, left, top, w, h, color, shape_type=1, radius=None):
        """Add a filled shape (1=rect, 5=rounded-rect, 9=oval). No border."""
        s = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(w), Inches(h))
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()
        if radius is not None:
            # set corner radius (EMU; 914400 = 1 inch; 45720 ≈ 0.05")
            try:
                s.adjustments[0] = radius
            except Exception:
                pass
        return s

    def _text(slide, left, top, w, h, txt, size=14, color=TEXT,
              bold=False, align=PP_ALIGN.LEFT, italic=False, wrap=True):
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.text = txt
        p.font.size    = Pt(size)
        p.font.color.rgb = color
        p.font.bold    = bold
        p.font.italic  = italic
        p.alignment    = align
        return tf

    def _footer(slide, slide_num, total):
        _text(slide, 0.45, 7.05, 5, 0.35, "SMARTAssist Hub",
              size=9, color=MUTED)
        _text(slide, 11.5, 7.05, 1.4, 0.35,
              f"{slide_num} / {total}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)

    def _section_header(slide, icon_char, title, accent, slide_num, total):
        """Icon pill + section title + footer."""
        pill = slide.shapes.add_shape(5, Inches(0.45), Inches(0.28), Inches(0.48), Inches(0.48))
        pill.fill.solid()
        pill.fill.fore_color.rgb = accent
        pill.line.fill.background()
        tf_pill = pill.text_frame
        tf_pill.paragraphs[0].text = icon_char
        tf_pill.paragraphs[0].font.size = Pt(18)
        tf_pill.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf_pill.paragraphs[0].font.color.rgb = TEXT

        _text(slide, 1.08, 0.30, 11, 0.50, title, size=22, color=TEXT, bold=True)
        _footer(slide, slide_num, total)

    def _bullet_row(slide, y, number, text, accent):
        """Numbered circle + body text row."""
        dot = slide.shapes.add_shape(9, Inches(0.45), Inches(y), Inches(0.34), Inches(0.34))
        dot.fill.solid()
        dot.fill.fore_color.rgb = accent
        dot.line.fill.background()
        tf_dot = dot.text_frame
        tf_dot.paragraphs[0].text = str(number)
        tf_dot.paragraphs[0].font.size = Pt(10)
        tf_dot.paragraphs[0].font.bold = True
        tf_dot.paragraphs[0].font.color.rgb = TEXT
        tf_dot.paragraphs[0].alignment = PP_ALIGN.CENTER
        _text(slide, 0.95, y - 0.02, 11.8, 0.45, text, size=13, color=TEXT)

    # ── Normalise multi-chart / multi-table inputs ────────────────────────────
    # New path: data["charts"] = [{"title":..., "image_b64":...}, ...]
    # Legacy path: single chart_image_b64 + data["chart"]["title"]
    raw_charts = data.get("charts", [])
    if raw_charts:
        charts_to_render = [
            {"title": c.get("title", ""), "image_bytes": _decode_chart_image(c.get("image_b64"))}
            for c in raw_charts if c.get("image_b64")
        ]
    else:
        legacy_bytes = _decode_chart_image(chart_image_b64)
        charts_to_render = (
            [{"title": (data.get("chart") or {}).get("title", ""), "image_bytes": legacy_bytes}]
            if legacy_bytes else []
        )

    # New path: data["tables"] = [{"headers":..., "rows":..., "label":...}, ...]
    # Legacy path: data["table"] / data["table2"]
    raw_tables = data.get("tables", [])
    if raw_tables:
        tables_to_render = [t for t in raw_tables if t.get("headers") and t.get("rows")]
    else:
        tables_to_render = []
        for key in ("table", "table2"):
            t = data.get(key)
            if t and t.get("headers") and t.get("rows"):
                tables_to_render.append(t)

    # ── Count slides for footer total ────────────────────────────────────────
    penemuan = data.get("penemuan", [])
    tafsiran = data.get("tafsiran", "")
    cadangan = data.get("cadangan", [])
    amaran   = data.get("amaran", [])

    total = 1  # title
    if penemuan:  total += 1
    if tafsiran:  total += 1
    total += len(charts_to_render)
    total += len(tables_to_render)
    if cadangan:  total += 1
    if amaran:    total += 1
    total += 1   # closing slide
    cur = 0

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE 1 — Title
    # ═══════════════════════════════════════════════════════════════════════
    cur += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)

    # Decorative circles (background accents — NOT spanning stripes)
    _shape(slide, 9.6, -2.2, 5.8, 5.8, BG2, shape_type=9)
    _shape(slide, 11.0, -0.6, 3.0, 3.0, RGBColor(0x1E, 0x3A, 0x8A), shape_type=9)
    _shape(slide, -1.0, 5.6,  3.4, 3.4, BG2, shape_type=9)

    # Brand label top-left
    _text(slide, 0.55, 0.48, 5, 0.38, "SMARTAssist Hub",
          size=11, color=C["gold"], bold=True)

    # Main title
    _text(slide, 0.55, 1.75, 9.5, 1.3, "ANALISIS DATA",
          size=48, color=TEXT, bold=True)

    # Accent dot
    _shape(slide, 0.55, 3.1, 0.13, 0.13, C["blue"], shape_type=9)

    # Subtitle / message
    msg = data.get("message", "Laporan Analisis Data")
    _text(slide, 0.75, 3.0, 9.2, 0.8, msg, size=16, color=SUBTLE)

    # Date bottom-left
    _text(slide, 0.55, 6.65, 8, 0.38,
          datetime.now().strftime("%d %B %Y"),
          size=10, color=MUTED)

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE — Penemuan (Key Findings)
    # ═══════════════════════════════════════════════════════════════════════
    if penemuan:
        cur += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(slide)
        _section_header(slide, "1", "Penemuan Utama", C["blue"], cur, total)

        # Content card
        _shape(slide, 0.45, 0.90, 12.4, 5.9, BG2, shape_type=1)

        y = 1.07
        for i, item in enumerate(penemuan[:10]):
            _bullet_row(slide, y, i + 1, item, C["blue"])
            y += 0.55
            if y > 6.55:
                break

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE — Tafsiran (Interpretation)
    # ═══════════════════════════════════════════════════════════════════════
    if tafsiran:
        cur += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(slide)
        _section_header(slide, "2", "Tafsiran & Analisis", C["cyan"], cur, total)

        # Content card
        _shape(slide, 0.45, 0.90, 12.4, 5.9, BG2, shape_type=1)

        # Decorative quote circle
        _shape(slide, 0.6, 1.0, 0.28, 0.28, C["cyan"], shape_type=9)

        _text(slide, 1.0, 1.05, 11.6, 5.55, tafsiran,
              size=14, color=TEXT, wrap=True)

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDES — Charts (one slide per chart)
    # ═══════════════════════════════════════════════════════════════════════
    chart_icons = ["3", "3b", "3c", "3d"]
    for c_idx, chart in enumerate(charts_to_render):
        if not chart["image_bytes"]:
            continue
        cur += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(slide)
        icon_label = str(3 + c_idx) if c_idx == 0 else f"C{c_idx + 1}"
        _section_header(slide, icon_label, f"Visualisasi Data {c_idx + 1}" if len(charts_to_render) > 1 else "Visualisasi Data", C["purple"], cur, total)

        chart_title = chart.get("title", "")
        if chart_title:
            _text(slide, 0.6, 0.9, 12, 0.38, chart_title,
                  size=12, color=SUBTLE, align=PP_ALIGN.CENTER)
            img_top = 1.35
        else:
            img_top = 1.0

        img_stream = io.BytesIO(chart["image_bytes"])
        slide.shapes.add_picture(img_stream,
                                 Inches(1.2), Inches(img_top),
                                 Inches(10.8), Inches(7.5 - img_top - 0.55))

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDES — Data Tables (one slide per table)
    # ═══════════════════════════════════════════════════════════════════════
    for t_idx, table_data in enumerate(tables_to_render):
        cur += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(slide)
        tbl_label = table_data.get("label", f"Jadual {t_idx + 1}" if len(tables_to_render) > 1 else "Jadual Data")
        icon_label = str(3 + len(charts_to_render) + t_idx)
        _section_header(slide, icon_label, tbl_label, C["green"], cur, total)

        headers = table_data["headers"]
        rows    = table_data["rows"]
        num_cols = len(headers)
        num_rows = min(len(rows), 14) + 1
        col_width = min(12.0 / num_cols, 2.8)
        tbl_width = col_width * num_cols

        tbl_shape = slide.shapes.add_table(
            num_rows, num_cols,
            Inches(0.55), Inches(1.0),
            Inches(tbl_width), Inches(0.38 * num_rows)
        )
        tbl = tbl_shape.table

        HDR_BG = RGBColor(0x06, 0x52, 0x4A)
        ROW_A  = BG2
        ROW_B  = RGBColor(0x16, 0x1E, 0x2E)

        for i, h in enumerate(headers):
            cell = tbl.cell(0, i)
            cell.text = str(h)
            for p in cell.text_frame.paragraphs:
                p.font.size      = Pt(10)
                p.font.color.rgb = TEXT
                p.font.bold      = True
            cell.fill.solid()
            cell.fill.fore_color.rgb = HDR_BG

        for r_idx, row in enumerate(rows[:14]):
            for c_idx2, val in enumerate(row):
                cell = tbl.cell(r_idx + 1, c_idx2)
                cell.text = str(val)
                for p in cell.text_frame.paragraphs:
                    p.font.size      = Pt(9)
                    p.font.color.rgb = TEXT
                cell.fill.solid()
                cell.fill.fore_color.rgb = ROW_A if r_idx % 2 == 0 else ROW_B

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE — Cadangan (Recommendations)
    # ═══════════════════════════════════════════════════════════════════════
    if cadangan:
        cur += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(slide)
        _section_header(slide, "5", "Cadangan & Tindakan", C["amber"], cur, total)

        _shape(slide, 0.45, 0.90, 12.4, 5.9, BG2, shape_type=1)

        y = 1.07
        for i, item in enumerate(cadangan[:10]):
            _bullet_row(slide, y, i + 1, item, C["amber"])
            y += 0.55
            if y > 6.55:
                break

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE — Amaran (Warnings)
    # ═══════════════════════════════════════════════════════════════════════
    if amaran:
        cur += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(slide)
        _section_header(slide, "!", "Amaran & Limitasi", C["red"], cur, total)

        _shape(slide, 0.45, 0.90, 12.4, 5.9, BG3, shape_type=1)

        y = 1.07
        for i, item in enumerate(amaran[:10]):
            _bullet_row(slide, y, i + 1, item, C["red"])
            y += 0.55
            if y > 6.55:
                break

    # ═══════════════════════════════════════════════════════════════════════
    # CLOSING SLIDE
    # ═══════════════════════════════════════════════════════════════════════
    cur += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)

    _shape(slide, 9.6, -2.2, 5.8, 5.8, BG2, shape_type=9)
    _shape(slide, -0.8, 5.2, 3.0, 3.0, BG2, shape_type=9)

    _text(slide, 1.5, 2.3, 10, 1.0, "Terima Kasih",
          size=40, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    _text(slide, 1.5, 3.4, 10, 0.5,
          "AI assistant, make your work easier.",
          size=15, color=SUBTLE, align=PP_ALIGN.CENTER)
    _text(slide, 1.5, 4.1, 10, 0.38, "SMARTAssist Hub",
          size=12, color=C["gold"], bold=True, align=PP_ALIGN.CENTER)
    _text(slide, 1.5, 4.55, 10, 0.32,
          datetime.now().strftime("%d %B %Y"),
          size=10, color=MUTED, align=PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def build_pdf(data: dict, chart_image_b64: str | None = None) -> bytes:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import cm
    from reportlab.lib.colors import HexColor
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.enums import TA_LEFT, TA_CENTER

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4, leftMargin=2*cm, rightMargin=2*cm, topMargin=2*cm, bottomMargin=2*cm)

    styles = {
        "title": ParagraphStyle("title", fontName="Helvetica-Bold", fontSize=20, textColor=HexColor("#3b82f6"), spaceAfter=12, alignment=TA_CENTER),
        "subtitle": ParagraphStyle("subtitle", fontName="Helvetica", fontSize=11, textColor=HexColor("#94a3b8"), spaceAfter=20, alignment=TA_CENTER),
        "heading": ParagraphStyle("heading", fontName="Helvetica-Bold", fontSize=14, textColor=HexColor("#3b82f6"), spaceBefore=16, spaceAfter=8),
        "body": ParagraphStyle("body", fontName="Helvetica", fontSize=11, textColor=HexColor("#1e293b"), spaceAfter=6, leading=16),
        "bullet": ParagraphStyle("bullet", fontName="Helvetica", fontSize=11, textColor=HexColor("#1e293b"), spaceAfter=4, leftIndent=20, leading=16),
        "warning": ParagraphStyle("warning", fontName="Helvetica", fontSize=11, textColor=HexColor("#b45309"), spaceAfter=4, leftIndent=20, leading=16),
    }

    elements = []

    elements.append(Paragraph("Analisis Data", styles["title"]))
    elements.append(Paragraph(f"SMARTAssist Hub — {datetime.now().strftime('%d %B %Y')}", styles["subtitle"]))

    msg = data.get("message", "")
    if msg:
        elements.append(Paragraph(msg, styles["body"]))
        elements.append(Spacer(1, 8))

    penemuan = data.get("penemuan", [])
    if penemuan:
        elements.append(Paragraph("\U0001F4CC Penemuan", styles["heading"]))
        for p in penemuan:
            elements.append(Paragraph(f"•  {p}", styles["bullet"]))

    tafsiran = data.get("tafsiran", "")
    if tafsiran:
        elements.append(Paragraph("\U0001F50D Tafsiran", styles["heading"]))
        elements.append(Paragraph(tafsiran, styles["body"]))

    # --- Chart image ---
    chart_bytes = _decode_chart_image(chart_image_b64)
    if chart_bytes:
        elements.append(Spacer(1, 8))
        chart_title = ""
        if data.get("chart") and data["chart"].get("title"):
            chart_title = data["chart"]["title"]
            elements.append(Paragraph(f"\U0001F4CA {chart_title}", styles["heading"]))
        else:
            elements.append(Paragraph("\U0001F4CA Carta", styles["heading"]))
        img_stream = io.BytesIO(chart_bytes)
        page_width = A4[0] - 4 * cm
        img = Image(img_stream, width=page_width, height=page_width * 0.5)
        elements.append(img)
        elements.append(Spacer(1, 8))

    table_data = data.get("table")
    if table_data and table_data.get("headers") and table_data.get("rows"):
        elements.append(Paragraph("\U0001F4CA Data", styles["heading"]))
        headers = table_data["headers"]
        rows = table_data["rows"][:20]
        tbl_data = [headers] + rows
        col_widths = [min(450 / len(headers), 120)] * len(headers)
        tbl = Table(tbl_data, colWidths=col_widths)
        tbl.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), HexColor("#334155")),
            ("TEXTCOLOR", (0, 0), (-1, 0), HexColor("#f1f5f9")),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), 9),
            ("TEXTCOLOR", (0, 1), (-1, -1), HexColor("#1e293b")),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [HexColor("#f8fafc"), HexColor("#e2e8f0")]),
            ("GRID", (0, 0), (-1, -1), 0.5, HexColor("#cbd5e1")),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ("TOPPADDING", (0, 0), (-1, -1), 4),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
        ]))
        elements.append(tbl)

    cadangan = data.get("cadangan", [])
    if cadangan:
        elements.append(Paragraph("\U0001F4A1 Cadangan", styles["heading"]))
        for c in cadangan:
            elements.append(Paragraph(f"•  {c}", styles["bullet"]))

    amaran = data.get("amaran", [])
    if amaran:
        elements.append(Paragraph("⚠️ Amaran", styles["heading"]))
        for a in amaran:
            elements.append(Paragraph(f"•  {a}", styles["warning"]))

    doc.build(elements)
    return buf.getvalue()


def build_xlsx(data: dict, df=None) -> bytes:
    """Export the analysis to an .xlsx workbook: summary, analysis table(s),
    an auto-computed pivot, and the full dataset."""
    import pandas as pd

    buf = io.BytesIO()
    with pd.ExcelWriter(buf, engine="openpyxl") as writer:
        # ── Summary sheet ──
        rows = []
        if data.get("message"):
            rows.append(["Ringkasan", data["message"]])
        for label, key in [("Penemuan", "penemuan"), ("Cadangan", "cadangan"), ("Amaran", "amaran")]:
            for item in (data.get(key) or []):
                rows.append([label, item])
        if data.get("tafsiran"):
            rows.append(["Tafsiran", data["tafsiran"]])
        summary_df = pd.DataFrame(rows or [["Ringkasan", "Analisis Data"]], columns=["Bahagian", "Butiran"])
        summary_df.to_excel(writer, sheet_name="Ringkasan", index=False)

        # ── Analysis table(s) ──
        for i, key in enumerate(["table", "table2"]):
            tbl = data.get(key)
            if tbl and tbl.get("headers") and tbl.get("rows"):
                tdf = pd.DataFrame(tbl["rows"], columns=tbl["headers"])
                tdf.to_excel(writer, sheet_name=("Jadual Analisis" if i == 0 else "Jadual Perbandingan"), index=False)

        # ── Auto pivot from the full dataset (first category × mean of numerics) ──
        if df is not None and not df.empty:
            cats = df.select_dtypes(include=["object", "category"]).columns.tolist()
            nums = df.select_dtypes(include=["number"]).columns.tolist()
            cats = [c for c in cats if 1 < df[c].nunique(dropna=True) <= 50]
            if cats and nums:
                try:
                    pivot = df.groupby(cats[0])[nums].mean().round(2)
                    pivot.to_excel(writer, sheet_name="Pivot")
                except Exception:
                    pass
            # ── Full dataset ──
            safe = df.copy()
            if safe.shape[0] > 100000:
                safe = safe.head(100000)
            safe.to_excel(writer, sheet_name="Data Penuh", index=False)

    return buf.getvalue()
